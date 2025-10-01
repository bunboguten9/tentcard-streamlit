# app.py （Streamlit・一意キー付け修正版）

import io
from copy import deepcopy
from math import sqrt, isfinite
from pathlib import Path
from typing import Optional, Tuple, List

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# ==============================
# 既定ファイル名（同じフォルダにテンプレが置いてある想定）
# ==============================
DEFAULT_TEMPLATE = Path("Tentcard_JP_format.pptx")
OUTPUT_FILENAME = "Tentcard_generated.pptx"

# ==============================
# 列名（Outputsシート）
# ==============================
COL_COMPANY_JP = "Company(JP)"
COL_FIRST      = "First Name"
COL_LAST       = "Last Name"

# ==============================
# 既定の書式（デフォルト）
# ==============================
GLOBAL_DEFAULT_FONT = "游ゴシック Regular"

COMPANY_SIZE_MAX = 44
COMPANY_SIZE_MIN = 20
COMPANY_LINE_SPACING_DEFAULT = 1.0
COMPANY_BOLD_DEFAULT = True

PERSON_SIZE_DEFAULT = 80
PERSON_LINE_SPACING_DEFAULT = 1.5
PERSON_BOLD_DEFAULT = True

# ==============================
# 二次近似モデル：日本語のみ会社名の「折返しなし最大文字数」 f(s)
# 実測点：44→17, 38→20, 31→24, 30→25, 28→27, 20→38 に最小二乗近似
# ==============================
A_COEF = 0.02952
B_COEF = -2.7438
C_COEF = 80.904

def f_allowed_chars(size_pt: float) -> float:
    return A_COEF * size_pt * size_pt + B_COEF * size_pt + C_COEF

# ==============================
# ユーティリティ
# ==============================
def norm(v) -> str:
    if pd.isna(v):
        return ""
    return str(v).strip()

def _has_textframe(shape) -> bool:
    return hasattr(shape, "text_frame") and shape.text_frame is not None

def find_tentcard_textboxes(slide) -> Tuple[Optional[object], Optional[object]]:
    """幅の大きい text_frame を持つ図形 上位2つを (上,下) で返す。"""
    tf_shapes = [s for s in slide.shapes if _has_textframe(s) and getattr(s, "width", None) is not None]
    if not tf_shapes:
        return None, None
    top2 = sorted(tf_shapes, key=lambda s: s.width, reverse=True)[:2]
    top2 = sorted(top2, key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)))
    if len(top2) == 2:
        return top2[0], top2[1]
    if len(top2) == 1:
        return top2[0], None
    return None, None

def _clear_paragraph_runs(p):
    while p.runs:
        r = p.runs[0]._r
        r.getparent().remove(r)

def _set_paragraph_text_with_style(
    p,
    text: str,
    *,
    font_name: str,
    size_pt: int,
    bold: bool,
    line_spacing_multiple: float,
):
    _clear_paragraph_runs(p)
    r = p.add_run()
    r.text = text

    f = r.font
    f.name = font_name
    f.size = Pt(size_pt)
    f.bold = bold

    if p.alignment is None:
        p.alignment = PP_ALIGN.CENTER
    p.line_spacing = float(line_spacing_multiple)

# ===== 日本語/英語名 判定 =====
def _is_japanese_char(ch: str) -> bool:
    code = ord(ch)
    if 0x3040 <= code <= 0x309F:  # ひらがな
        return True
    if 0x30A0 <= code <= 0x30FF or 0x31F0 <= code <= 0x31FF:  # カタカナ系
        return True
    if 0x4E00 <= code <= 0x9FFF:  # 漢字
        return True
    if 0xFF01 <= code <= 0xFF60 or 0xFFE0 <= code <= 0xFFE6:  # 全角記号
        return True
    if 0x3000 <= code <= 0x303F:  # CJK記号（・「」等）
        return True
    return False

def is_all_japanese(s: str) -> bool:
    s2 = s.replace(" ", "").replace("\u3000", "")
    if not s2:
        return False
    return all(_is_japanese_char(ch) for ch in s2)

def is_english_name(first: str, last: str) -> bool:
    s = (first or "") + (last or "")
    for ch in s:
        o = ord(ch)
        if (65 <= o <= 90) or (97 <= o <= 122):  # A-Z, a-z
            return True
    return False

def effective_len_jp(s: str) -> int:
    return len(s.replace(" ", "").replace("\u3000", ""))

def compute_company_font_size(company_text: str, fixed_override: Optional[int]) -> int:
    """
    フォントサイズの決定：
      - 固定サイズが指定されていればそれを採用（クリップ有り）
      - 未指定なら：日本語のみ → 放物線モデル逆解（20〜44）、英字混在 → 44固定
    """
    if fixed_override is not None:
        return max(COMPANY_SIZE_MIN, min(COMPANY_SIZE_MAX, int(fixed_override)))

    if not company_text:
        return COMPANY_SIZE_MAX
    if not is_all_japanese(company_text):
        return COMPANY_SIZE_MAX

    n = effective_len_jp(company_text)
    if n <= 0:
        return COMPANY_SIZE_MAX

    if n <= f_allowed_chars(COMPANY_SIZE_MAX) + 1e-6:
        return COMPANY_SIZE_MAX
    if n >= f_allowed_chars(COMPANY_SIZE_MIN) - 1e-6:
        return COMPANY_SIZE_MIN

    # 逆解（小さい方の解）
    a, b, c = A_COEF, B_COEF, C_COEF - n
    D = b*b - 4*a*c
    if D < 0:
        return COMPANY_SIZE_MIN
    s1 = (-b - sqrt(D)) / (2*a)
    s2 = (-b + sqrt(D)) / (2*a)
    s_candidate = s1 if s1 <= s2 else s2
    if not isfinite(s_candidate):
        return COMPANY_SIZE_MIN

    size = int(s_candidate)  # 安全側に切り捨て
    size = max(COMPANY_SIZE_MIN, min(COMPANY_SIZE_MAX, size))
    return size

# ===== 会社名の「株式会社」処理 =====
def apply_kabushikigaisha(company_text: str, include_kabu: bool) -> str:
    if not company_text:
        return company_text
    if include_kabu:
        return company_text  # 現状維持
    # 先頭の「株式会社」を1回だけ除去（前方空白は温存）
    stripped = company_text.lstrip()
    if stripped.startswith("株式会社"):
        left_pad = company_text[:len(company_text) - len(stripped)]
        return left_pad + stripped[len("株式会社"):]
    return company_text

# ===== スライドへの書き込み（詳細設定を考慮） =====
def set_company_and_name(
    slide,
    company_jp: str,
    first_name: str,
    last_name: str,
    *,
    # グローバル
    global_font: str,
    # 会社名：基本
    include_kabu: bool,
    # 会社名：詳細設定（None/空は既定）
    company_font_override: Optional[str],
    company_fixed_size: Optional[int],
    company_line_spacing_override: Optional[float],
    company_bold_mode: str,  # "既定","ON","OFF"
    # 氏名：基本
    add_sama: bool,
    add_sama_for_english: bool,
    # 氏名：詳細設定
    person_font_override: Optional[str],
    person_size_override: Optional[int],
    person_line_spacing_override: Optional[float],
    person_bold_mode: str,   # "既定","ON","OFF"
):
    # フォント名の決定
    base_font = (global_font.strip() if global_font.strip() else GLOBAL_DEFAULT_FONT)
    company_font = (company_font_override.strip() if company_font_override and company_font_override.strip() else base_font)
    person_font  = (person_font_override.strip()  if person_font_override  and person_font_override.strip()  else base_font)

    # 会社名の表示文字列
    company_disp = apply_kabushikigaisha(company_jp, include_kabu)

    # 氏名表示＆様
    full_name = f"{last_name} {first_name}".strip()
    if add_sama:
        if add_sama_for_english or not is_english_name(first_name, last_name):
            full_name = f"{full_name} 様"

    # 書式（太字）決定
    if company_bold_mode == "ON":
        company_bold = True
    elif company_bold_mode == "OFF":
        company_bold = False
    else:
        company_bold = COMPANY_BOLD_DEFAULT

    if person_bold_mode == "ON":
        person_bold = True
    elif person_bold_mode == "OFF":
        person_bold = False
    else:
        person_bold = PERSON_BOLD_DEFAULT

    # 行間の決定
    company_ls = company_line_spacing_override if company_line_spacing_override else COMPANY_LINE_SPACING_DEFAULT
    person_ls  = person_line_spacing_override  if person_line_spacing_override  else PERSON_LINE_SPACING_DEFAULT

    # テキスト枠
    top_box, bottom_box = find_tentcard_textboxes(slide)

    def _fill(box):
        if box is None or box.text_frame is None:
            return
        tf = box.text_frame
        tf.clear()

        # 段落1：会社名（サイズは自動 or 固定）
        company_size = compute_company_font_size(company_disp, company_fixed_size)
        p1 = tf.paragraphs[0]
        _set_paragraph_text_with_style(
            p1, company_disp,
            font_name=company_font,
            size_pt=company_size,
            bold=company_bold,
            line_spacing_multiple=company_ls,
        )

        # 段落2：氏名（サイズは既定80 or 上書き）
        person_size = int(person_size_override) if person_size_override else PERSON_SIZE_DEFAULT
        p2 = tf.add_paragraph()
        _set_paragraph_text_with_style(
            p2, full_name,
            font_name=person_font,
            size_pt=person_size,
            bold=person_bold,
            line_spacing_multiple=person_ls,
        )

    _fill(top_box)
    _fill(bottom_box)

# ==============================
# プレビュー（文字なし：枠の位置とサイズを図示）
# ==============================
def render_layout_preview(prs: Presentation):
    if not prs.slides:
        st.info("テンプレートにスライドがありません。")
        return
    slide = prs.slides[0]
    top_box, bottom_box = find_tentcard_textboxes(slide)

    W = prs.slide_width
    H = prs.slide_height

    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_title("レイアウトプレビュー（文字なし）", fontsize=12)
    ax.set_xlabel("EMU (x)")
    ax.set_ylabel("EMU (y)")

    # スライド枠
    ax.add_patch(plt.Rectangle((0, 0), W, H, fill=False, linewidth=2))

    # テキストボックス（上下）
    def draw_box(shp, color="C1"):
        if shp is None:
            return
        x, y, w, h = shp.left, shp.top, shp.width, shp.height
        ax.add_patch(plt.Rectangle((x, y), w, h, fill=False, linewidth=2, linestyle="--", color=color))

    draw_box(top_box, "C2")
    draw_box(bottom_box, "C3")

    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)  # PowerPoint座標は上が0なので上下反転
    st.pyplot(fig)

# ==============================
# Streamlit UI
# ==============================
st.set_page_config(page_title="Tentcard Generator", layout="centered")
st.header("テントカード自動生成")

# --- ファイル入力 ---
excel_file = st.file_uploader("Excel（Outputsシート）をドラッグ＆ドロップ（.xlsx）", type=["xlsx"], key="file_excel")
template_file = st.file_uploader("PPTXテンプレ（省略可：未指定なら既定を使用）", type=["pptx"], key="file_template")

# --- 全体フォント ---
global_font = st.text_input('全体のフォント名（空欄OK / 空欄時は「游ゴシック Regular」）', value="", key="global_font")

# --- 会社名（基本） ---
st.subheader("会社名（段落1）")
include_kabu = st.checkbox("会社名に「株式会社」を含める", value=True, key="company_include_kabu")

with st.expander("会社名の詳細設定", expanded=False):
    company_font_override = st.text_input("会社名のフォント名（空欄=全体フォントを使用）", value="", key="company_font_override")
    use_company_fixed = st.checkbox("会社名のフォントサイズを固定する", value=False, key="company_use_fixed")
    company_fixed_size = None
    if use_company_fixed:
        company_fixed_size = st.number_input("固定フォントサイズ（pt）", min_value=10, max_value=200, value=44, step=1, key="company_fixed_size")
    company_bold_mode = st.selectbox("太字", options=["既定", "ON", "OFF"], index=0, key="company_bold_mode")
    company_line_spacing_txt = st.text_input("行間（倍率 / 空欄=既定1.0）", value="", key="company_line_spacing_txt")
    company_line_spacing_override = float(company_line_spacing_txt) if company_line_spacing_txt.strip() else None

# --- 氏名（基本） ---
st.subheader("氏名（段落2）")
add_sama = st.checkbox("氏名に「様」を付与", value=True, key="person_add_sama")
add_sama_for_english = st.checkbox("英語名にも「様」を付与", value=False, key="person_add_sama_en")

with st.expander("氏名の詳細設定", expanded=False):
    person_font_override = st.text_input("氏名のフォント名（空欄=全体フォントを使用）", value="", key="person_font_override")
    use_person_fixed = st.checkbox("氏名のフォントサイズを変更する", value=False, key="person_use_fixed")
    person_size_override = None
    if use_person_fixed:
        person_size_override = st.number_input("氏名フォントサイズ（pt）", min_value=10, max_value=200, value=PERSON_SIZE_DEFAULT, step=1, key="person_size_override")
    person_bold_mode = st.selectbox("太字", options=["既定", "ON", "OFF"], index=0, key="person_bold_mode")
    person_line_spacing_txt = st.text_input("行間（倍率 / 空欄=既定1.5）", value="", key="person_line_spacing_txt")
    person_line_spacing_override = float(person_line_spacing_txt) if person_line_spacing_txt.strip() else None

st.divider()

# --- プレビュー ---
if st.button("レイアウトプレビュー（文字なし）", key="btn_preview"):
    try:
        if template_file is not None:
            prs_prev = Presentation(io.BytesIO(template_file.read()))
        else:
            if not DEFAULT_TEMPLATE.exists():
                st.error("既定テンプレートが見つかりませんでした。テンプレPPTXをアップロードしてください。")
            else:
                prs_prev = Presentation(str(DEFAULT_TEMPLATE))
        render_layout_preview(prs_prev)
    except Exception as e:
        st.error(f"プレビューに失敗しました: {e}")

st.divider()

# --- 生成 ---
generate = st.button("PPTXを生成（1人=1スライド集約）", key="btn_generate")

if generate:
    # Excelチェック
    if excel_file is None:
        st.error("Excelファイルを指定してください。")
        st.stop()

    # Excel読込（Outputs固定）
    try:
        df = pd.read_excel(excel_file, sheet_name="Outputs", dtype=str)
    except Exception as e:
        st.error(f"Excelの読み込みに失敗しました: {e}")
        st.stop()

    # 列チェック
    required = [COL_COMPANY_JP, COL_FIRST, COL_LAST]
    missing = [c for c in required if c not in df.columns]
    if missing:
        st.error(f"必要な列が不足しています: {missing}")
        st.stop()

    # テンプレ読込
    try:
        if template_file is not None:
            prs = Presentation(io.BytesIO(template_file.read()))
        else:
            if not DEFAULT_TEMPLATE.exists():
                st.error("既定テンプレートが見つかりません。テンプレPPTXをアップロードしてください。")
                st.stop()
            prs = Presentation(str(DEFAULT_TEMPLATE))
    except Exception as e:
        st.error(f"PPTXテンプレートの読み込みに失敗しました: {e}")
        st.stop()

    if len(prs.slides) == 0:
        st.error("テンプレートにスライドがありません。")
        st.stop()

    # レコード取得（Excel順のまま）
    records: List[Tuple[str, str, str]] = []
    for _, row in df.iterrows():
        company = norm(row.get(COL_COMPANY_JP))
        first   = norm(row.get(COL_FIRST))
        last    = norm(row.get(COL_LAST))
        if not company or not first or not last:
            continue
        records.append((company, first, last))

    if not records:
        st.warning("有効な行がありませんでした。")
        st.stop()

    # 1件目をテンプレ1枚目へ
    template_slide = prs.slides[0]
    set_company_and_name(
        template_slide,
        records[0][0], records[0][1], records[0][2],
        global_font=global_font,
        include_kabu=include_kabu,
        company_font_override=company_font_override,
        company_fixed_size=int(company_fixed_size) if company_fixed_size is not None else None,
        company_line_spacing_override=company_line_spacing_override,
        company_bold_mode=company_bold_mode,
        add_sama=add_sama,
        add_sama_for_english=add_sama_for_english,
        person_font_override=person_font_override,
        person_size_override=int(person_size_override) if person_size_override is not None else None,
        person_line_spacing_override=person_line_spacing_override,
        person_bold_mode=person_bold_mode,
    )

    # スライド複製
    def duplicate_slide(prs: Presentation, src_slide):
        try:
            blank = prs.slide_layouts[6]
        except IndexError:
            blank = prs.slide_layouts[len(prs.slide_layouts) - 1]
        new_slide = prs.slides.add_slide(blank)
        spTree = new_slide.shapes._spTree
        src_spTree = src_slide.shapes._spTree
        for el in src_spTree:
            if el.tag.endswith("extLst"):
                continue
            new_el = deepcopy(el)
            spTree.insert_element_before(new_el, "p:extLst")
        return new_slide

    for company, first, last in records[1:]:
        sld = duplicate_slide(prs, template_slide)
        set_company_and_name(
            sld,
            company, first, last,
            global_font=global_font,
            include_kabu=include_kabu,
            company_font_override=company_font_override,
            company_fixed_size=int(company_fixed_size) if company_fixed_size is not None else None,
            company_line_spacing_override=company_line_spacing_override,
            company_bold_mode=company_bold_mode,
            add_sama=add_sama,
            add_sama_for_english=add_sama_for_english,
            person_font_override=person_font_override,
            person_size_override=int(person_size_override) if person_size_override is not None else None,
            person_line_spacing_override=person_line_spacing_override,
            person_bold_mode=person_bold_mode,
        )

    # バイナリ保存→ダウンロード
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    st.success(f"生成完了: {OUTPUT_FILENAME}")
    st.download_button(
        label="生成ファイルをダウンロード",
        data=out,
        file_name=OUTPUT_FILENAME,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="btn_download",
    )
